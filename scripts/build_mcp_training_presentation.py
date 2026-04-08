from __future__ import annotations

import argparse
import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_TEMPLATE = Path("/Users/pawansi/Downloads/Cisco_PowerPoint_Template_LIGHT_04-01-2026.potx")
DEFAULT_OUTPUT = Path("Plan_Provision_Validate_MCP_Enabled_Network_Automation.pptx")

BLUE = RGBColor(0, 84, 159)
TEAL = RGBColor(0, 170, 190)
SKY = RGBColor(223, 242, 248)
LIGHT = RGBColor(242, 247, 250)
MID = RGBColor(210, 221, 231)
DARK = RGBColor(32, 45, 64)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(44, 134, 92)
AMBER = RGBColor(242, 176, 52)


def convert_template_path(template_path: Path) -> tuple[Path, tempfile.TemporaryDirectory[str] | None]:
    if template_path.suffix.lower() != ".potx":
        return template_path, None

    temp_dir = tempfile.TemporaryDirectory(prefix="cisco_template_")
    converted = Path(temp_dir.name) / f"{template_path.stem}.pptx"
    shutil.copyfile(template_path, converted)
    with zipfile.ZipFile(converted, "r") as source:
        package = {name: source.read(name) for name in source.namelist()}

    content_types = package["[Content_Types].xml"].decode("utf-8")
    content_types = content_types.replace(
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    )

    with zipfile.ZipFile(converted, "w", zipfile.ZIP_DEFLATED) as target:
        for name, data in package.items():
            if name == "[Content_Types].xml":
                data = content_types.encode("utf-8")
            target.writestr(name, data)

    return converted, temp_dir


def remove_all_slides(prs: Presentation):
    for index in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[index]
        relationship_id = slide_id.rId
        prs.part.drop_rel(relationship_id)
        del prs.slides._sldIdLst[index]


def get_layout(prs: Presentation, layout_name: str):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    raise ValueError(f"Layout '{layout_name}' not found in template.")


def title_placeholder(slide):
    if slide.shapes.title is not None:
        return slide.shapes.title
    for placeholder in slide.placeholders:
        ph_type = placeholder.placeholder_format.type
        if ph_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
            return placeholder
    raise ValueError("Slide has no title placeholder.")


def content_placeholders(slide):
    placeholders = []
    for placeholder in slide.placeholders:
        ph_type = placeholder.placeholder_format.type
        if ph_type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.FOOTER,
            PP_PLACEHOLDER.SLIDE_NUMBER,
            PP_PLACEHOLDER.DATE,
        }:
            continue
        placeholders.append(placeholder)
    return sorted(placeholders, key=lambda item: (item.top, item.left))


def clear_text_frame(text_frame):
    text_frame.clear()
    text_frame.word_wrap = True
    return text_frame


def set_paragraph_runs(paragraph, text: str, font_size: int = 18, bold: bool = False, color: RGBColor | None = None):
    paragraph.text = ""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def fill_placeholder_lines(
    placeholder,
    lines: list[str],
    *,
    font_size: int = 18,
    first_line_bold: bool = False,
    color: RGBColor | None = None,
    center: bool = False,
):
    text_frame = clear_text_frame(placeholder.text_frame)
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        if center:
            paragraph.alignment = PP_ALIGN.CENTER
        set_paragraph_runs(
            paragraph,
            line,
            font_size=font_size,
            bold=first_line_bold and index == 0,
            color=color,
        )


def set_title(slide, text: str):
    title = title_placeholder(slide)
    text_frame = clear_text_frame(title.text_frame)
    set_paragraph_runs(text_frame.paragraphs[0], text, font_size=28, bold=True)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    font_size: int = 16,
    bold_first: bool = False,
    color: RGBColor = DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = align
        set_paragraph_runs(paragraph, line, font_size=font_size, bold=bold_first and index == 0, color=color)
    return textbox


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body_lines: list[str],
    *,
    fill_color: RGBColor = WHITE,
    line_color: RGBColor = MID,
    title_color: RGBColor = BLUE,
    body_color: RGBColor = DARK,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color

    add_textbox(slide, left + 0.16, top + 0.14, width - 0.32, 0.35, [title], font_size=18, bold_first=True, color=title_color)
    add_textbox(slide, left + 0.16, top + 0.55, width - 0.32, height - 0.7, body_lines, font_size=12, color=body_color)
    return shape


def add_step_card(slide, left: float, top: float, step: str, title: str, body: str, fill_color: RGBColor):
    add_card(
        slide,
        left,
        top,
        2.2,
        1.45,
        f"{step}. {title}",
        [body],
        fill_color=fill_color,
        line_color=fill_color,
        title_color=DARK,
    )


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = BLUE):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    connector.line.end_arrowhead = True
    return connector


def add_table_slide(slide, rows: int, cols: int, data: list[list[str]], col_widths: list[float]):
    table_placeholder = None
    for placeholder in content_placeholders(slide):
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.TABLE:
            table_placeholder = placeholder
            break
    if table_placeholder is None:
        raise ValueError("Slide has no table placeholder.")

    graphic_frame = table_placeholder.insert_table(rows, cols)
    table = graphic_frame.table
    for col_idx, col_width in enumerate(col_widths):
        table.columns[col_idx].width = Inches(col_width)

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[row_idx][col_idx]
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_idx % 2 else LIGHT
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            if row_idx == 0:
                set_paragraph_runs(paragraph, data[row_idx][col_idx], font_size=14, bold=True, color=WHITE)
            else:
                set_paragraph_runs(paragraph, data[row_idx][col_idx], font_size=12, color=DARK)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title Slide 2, Two Speakers"))
    set_title(slide, "Plan, Provision, Validate")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(placeholders[0], ["MCP-Enabled Network Automation for Catalyst Center teams"], font_size=20)
    fill_placeholder_lines(
        placeholders[1],
        ["Audience", "Catalyst Center operators and architects"],
        font_size=18,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[2],
        ["Starting point", "Strong Catalyst Center knowledge, little MCP or AI background"],
        font_size=18,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[3],
        ["Outcome", "A practical model for exposing trusted automation to AI clients"],
        font_size=16,
        first_line_bold=True,
    )


def add_agenda_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, Table 1"))
    set_title(slide, "Session Agenda")
    fill_placeholder_lines(content_placeholders(slide)[0], ["45 minutes, built for Catalyst Center practitioners"], font_size=18)
    agenda_rows = [
        ["Section", "Focus", "Minutes"],
        ["1. Why this matters", "What MCP solves and why existing automation is not enough by itself", "8"],
        ["2. Plan", "Exporting current Catalyst Center intent into reusable YAML", "8"],
        ["3. Provision", "Using safe tool contracts over workflow_manager modules", "10"],
        ["4. Validate and operate", "Tasks, progress, guardrails, and production deployment", "12"],
        ["5. Customer next steps", "Use cases, adoption path, and Q&A", "7"],
    ]
    add_table_slide(slide, len(agenda_rows), 3, agenda_rows, [2.8, 8.0, 1.2])


def add_segue(prs: Presentation, title: str):
    slide = prs.slides.add_slide(get_layout(prs, "Segue 1"))
    set_title(slide, title)


def add_existing_assets_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, 2 Columns"))
    set_title(slide, "You already own most of the automation")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(
        placeholders[0],
        ["The gap is not Catalyst Center capability. The gap is a safe contract between people, AI, and your existing automation."],
        font_size=16,
    )
    fill_placeholder_lines(
        placeholders[1],
        [
            "What teams already have",
            "Catalyst Center design, templates, assurance, and provisioning workflows",
            "Ansible collections, playbooks, Python helpers, and approved change processes",
            "Site standards, network settings baselines, and operating runbooks",
        ],
        font_size=15,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[2],
        [
            "What is still hard",
            "Choosing the right automation path from a natural-language request",
            "Shielding AI clients from nested module payloads and unsafe raw inputs",
            "Handling approvals, progress, retries, and audit for long-running changes",
        ],
        font_size=15,
        first_line_bold=True,
    )


def add_mcp_plain_english_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only 1"))
    set_title(slide, "MCP in Plain English")
    add_textbox(
        slide,
        0.35,
        1.0,
        12.2,
        0.5,
        ["Think of MCP as the tool contract that lets an AI client use your real automation without inventing API calls."],
        font_size=18,
    )
    add_card(slide, 0.45, 1.8, 2.9, 2.1, "1. User request", ["The operator asks for an outcome, such as exporting site intent or provisioning a building."], fill_color=LIGHT)
    add_card(slide, 3.55, 1.8, 2.9, 2.1, "2. MCP tool", ["The server advertises approved tools with typed inputs and clear descriptions."], fill_color=LIGHT)
    add_card(slide, 6.65, 1.8, 2.9, 2.1, "3. Real execution", ["The MCP server runs Ansible, SDK, or Terraform logic that already exists in your environment."], fill_color=LIGHT)
    add_card(slide, 9.75, 1.8, 2.9, 2.1, "4. Controlled result", ["The client gets task status, progress, artifacts, and final outputs instead of a guess."], fill_color=LIGHT)
    add_arrow(slide, 3.2, 2.85, 3.5, 2.85)
    add_arrow(slide, 6.3, 2.85, 6.6, 2.85)
    add_arrow(slide, 9.4, 2.85, 9.7, 2.85)
    add_card(
        slide,
        2.2,
        4.45,
        8.95,
        1.25,
        "Why this matters for Catalyst Center",
        [
            "Your approved automation stays in control. MCP adds a natural-language entry point, schema discipline, task tracking, and client interoperability."
        ],
        fill_color=SKY,
        line_color=SKY,
    )


def add_architecture_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only 1"))
    set_title(slide, "Catalyst Center MCP Server Architecture")
    add_textbox(
        slide,
        0.35,
        1.0,
        12.0,
        0.55,
        ["The AI does not talk directly to Catalyst Center APIs or credentials. The MCP server is the control point."],
        font_size=18,
    )
    add_card(slide, 0.45, 2.0, 2.05, 1.25, "Operator", ["Asks for an outcome in plain language"], fill_color=LIGHT)
    add_card(slide, 2.75, 2.0, 2.05, 1.25, "AI client", ["Chooses from advertised MCP tools"], fill_color=LIGHT)
    add_card(slide, 5.05, 1.75, 2.35, 1.75, "Catalyst Center MCP server", ["Flat tool schemas", "Translation layer", "Task primitives"], fill_color=SKY, line_color=SKY)
    add_card(slide, 7.7, 2.0, 2.15, 1.25, "ansible-runner", ["Runs workflow_manager and config_generator modules"], fill_color=LIGHT)
    add_card(slide, 10.1, 2.0, 2.35, 1.25, "Catalyst Center", ["Authoritative source of network intent and state"], fill_color=LIGHT)
    add_arrow(slide, 2.35, 2.62, 2.7, 2.62)
    add_arrow(slide, 4.65, 2.62, 5.0, 2.62)
    add_arrow(slide, 7.3, 2.62, 7.65, 2.62)
    add_arrow(slide, 9.75, 2.62, 10.05, 2.62)

    add_card(slide, 4.2, 4.3, 2.0, 1.15, "OAuth / SSO", ["Maps user identity to a tenant or policy boundary"], fill_color=LIGHT)
    add_card(slide, 6.45, 4.3, 2.0, 1.15, "Secrets", ["Catalyst Center credentials stay on the server side"], fill_color=LIGHT)
    add_card(slide, 8.7, 4.3, 2.1, 1.15, "Redis", ["Stores task state, progress, and artifacts index"], fill_color=LIGHT)

    add_arrow(slide, 6.2, 3.5, 5.9, 4.25, color=TEAL)
    add_arrow(slide, 6.35, 3.5, 7.05, 4.25, color=TEAL)
    add_arrow(slide, 6.55, 3.5, 9.0, 4.25, color=TEAL)


def add_plan_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only 1"))
    set_title(slide, "Plan Starts with Brownfield Export")
    add_textbox(
        slide,
        0.35,
        1.0,
        12.0,
        0.5,
        ["For Catalyst Center teams, the safest first AI use case is read-only export: pull current intent, review it, then decide what to change."],
        font_size=18,
    )
    add_step_card(slide, 0.45, 2.0, "1", "Generate", "Use playbook_config_generator modules to export current state into YAML.", fill_color=LIGHT)
    add_step_card(slide, 2.95, 2.0, "2", "Review", "Inspect the generated file as a normal Catalyst Center design artifact.", fill_color=LIGHT)
    add_step_card(slide, 5.45, 2.0, "3", "Edit", "Adjust the YAML only where a human has approved the intended change.", fill_color=LIGHT)
    add_step_card(slide, 7.95, 2.0, "4", "Apply", "Reapply through workflow_manager modules in merged mode.", fill_color=SKY)

    add_card(
        slide,
        0.65,
        4.2,
        4.8,
        1.45,
        "Useful starting generators",
        [
            "Site hierarchy",
            "Inventory",
            "Network settings",
            "Wireless design where the cluster exposes the APIs",
        ],
        fill_color=WHITE,
    )

    code_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.0),
        Inches(4.0),
        Inches(6.2),
        Inches(1.9),
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = DARK
    code_box.line.color.rgb = DARK
    code_lines = [
        "config:",
        "- network_management_details:",
        "  - site_name: Global/USA/SAN JOSE",
        "    settings:",
        "      timezone: GMT",
        "      syslog_server:",
        "        configure_dnac_ip: true",
    ]
    code_text = slide.shapes.add_textbox(Inches(6.18), Inches(4.18), Inches(5.85), Inches(1.55))
    tf = code_text.text_frame
    tf.word_wrap = False
    for index, line in enumerate(code_lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        set_paragraph_runs(paragraph, line, font_size=13, color=WHITE)
        paragraph.font.name = "Menlo"


def add_prompt_patterns_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, 2 Columns with Bullets"))
    set_title(slide, "What Good AI Requests Look Like")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(
        placeholders[0],
        [
            "Prompt patterns that work well",
            "Generate the site config for Global/USA/SAN JOSE as YAML",
            "Reapply the exported network settings in merged mode",
            "Create a sample building with two floors under SAN JOSE",
            "Show task progress and final result for the change",
        ],
        font_size=15,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[1],
        [
            "Why these prompts work",
            "They describe outcomes, not raw API calls",
            "They rely on typed tool schemas instead of nested module dictionaries",
            "They create a clean handoff to task IDs, artifacts, and approvals",
            "They keep Catalyst Center operators in the decision loop",
        ],
        font_size=15,
        first_line_bold=True,
    )


def add_tool_model_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, 4 Columns"))
    set_title(slide, "Tool Model for a Catalyst Center Audience")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(
        placeholders[0],
        ["The server can expose dozens of modules, but operators should learn three patterns first."],
        font_size=16,
    )
    fill_placeholder_lines(
        placeholders[1],
        ["Specialized tools", "Best for common tasks such as site provisioning, template deploy, inventory, and assurance."],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[2],
        ["Generic workflow tools", "run_*_workflow_manager wrappers cover the breadth of the Ansible collection when needed."],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[3],
        ["Config generators", "generate_*_config wrappers are the safest starting point for brownfield discovery and export."],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[4],
        ["Task primitives", "Every operation returns a taskId immediately so long-running changes are observable and pollable."],
        font_size=14,
        first_line_bold=True,
    )


def add_ppv_flow_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only 1"))
    set_title(slide, "Plan, Provision, Validate as One Operating Flow")
    add_textbox(
        slide,
        0.35,
        1.0,
        12.2,
        0.5,
        ["The power of MCP is not one tool. It is how an AI client chains trusted tools without losing control or auditability."],
        font_size=18,
    )
    steps = [
        ("Plan", "Pull current intent and generate reapply-safe YAML"),
        ("Approve", "Human review decides what is allowed to change"),
        ("Provision", "Run workflow_manager modules through ansible-runner"),
        ("Validate", "Use assurance, exports, and post-checks to confirm outcome"),
        ("Record", "Store tasks, artifacts, and change evidence"),
    ]
    left = 0.45
    for idx, (title, body) in enumerate(steps):
        fill = LIGHT if idx != 2 else SKY
        add_card(slide, left, 2.2, 2.35, 1.55, title, [body], fill_color=fill, line_color=fill)
        if idx < len(steps) - 1:
            add_arrow(slide, left + 2.35, 2.98, left + 2.55, 2.98)
        left += 2.55

    add_card(
        slide,
        2.1,
        4.55,
        9.1,
        1.2,
        "Important principle",
        ["The AI should orchestrate approved building blocks. It should not invent state models or bypass Catalyst Center governance."],
        fill_color=WHITE,
    )


def add_task_lifecycle_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only 1"))
    set_title(slide, "Long-Running Task Lifecycle")
    add_textbox(
        slide,
        0.35,
        1.0,
        12.2,
        0.5,
        ["Network changes take time. The MCP task model keeps the client responsive while ansible-runner does the work."],
        font_size=18,
    )
    add_step_card(slide, 0.45, 2.0, "1", "Submit", "Client calls a tool and gets taskId + submitted status immediately.", fill_color=LIGHT)
    add_step_card(slide, 2.95, 2.0, "2", "Execute", "ansible-runner starts the selected workflow_manager or config generator.", fill_color=LIGHT)
    add_step_card(slide, 5.45, 2.0, "3", "Track", "Runner events update progress, status messages, and artifact references.", fill_color=SKY)
    add_step_card(slide, 7.95, 2.0, "4", "Poll", "Client checks /tasks/get/{taskId} or receives progress notifications.", fill_color=LIGHT)
    add_step_card(slide, 10.45, 2.0, "5", "Complete", "Final result, generated files, or failure details are returned predictably.", fill_color=LIGHT)
    for start in [2.65, 5.15, 7.65, 10.15]:
        add_arrow(slide, start, 2.72, start + 0.2, 2.72)

    add_card(
        slide,
        1.1,
        4.35,
        11.1,
        1.25,
        "Why operators care",
        ["This is the difference between a chat demo and a production tool. You can see status, preserve evidence, and recover cleanly from failure."],
        fill_color=WHITE,
    )


def add_security_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, 3 Columns"))
    set_title(slide, "Security and Governance Model")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(
        placeholders[0],
        ["The recommended pattern is server-side secrets, tenant-aware access, and explicit approvals for risky operations."],
        font_size=16,
    )
    fill_placeholder_lines(
        placeholders[1],
        [
            "Identity and access",
            "Use enterprise SSO and OAuth to map each user to a tenant or policy boundary",
            "Authorize task reads so one tenant cannot inspect another tenant's changes",
        ],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[2],
        [
            "Secrets and approvals",
            "Keep Catalyst Center credentials on the MCP server, not in prompts",
            "Require confirmation metadata for destructive tools such as delete or reprovision actions",
        ],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[3],
        [
            "Operational controls",
            "Run over HTTPS with reverse proxy support, Redis-backed task state, and artifact retention rules",
            "Audit task inputs, progress events, and final outcomes for change records",
        ],
        font_size=14,
        first_line_bold=True,
    )


def add_deployment_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, 2 Columns"))
    set_title(slide, "Production Deployment Pattern")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(
        placeholders[0],
        ["A Linux deployment behind HTTPS gives customers a clean way to operate the server at scale."],
        font_size=16,
    )
    fill_placeholder_lines(
        placeholders[1],
        [
            "Core stack",
            "FastMCP + FastAPI for the tool endpoint",
            "ansible-runner as the only execution path",
            "Redis for task state and artifact indexing",
            "NGINX or a cloud load balancer for HTTPS and proxy buffering control",
        ],
        font_size=15,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[2],
        [
            "What this enables",
            "Stateless horizontal scale for the API tier",
            "MCP clients over HTTP/SSE",
            "Tenant-scoped credential routing",
            "Systemd or container-based deployment options",
        ],
        font_size=15,
        first_line_bold=True,
    )


def add_demo_storyboard_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, Table 1"))
    set_title(slide, "Suggested Demo Storyboard")
    fill_placeholder_lines(content_placeholders(slide)[0], ["Show one read-only export, one approved change, and one validation check."], font_size=17)
    rows = [
        ["Demo step", "Tool pattern", "What the audience should notice"],
        ["Export site or network settings", "generate_*_config", "Brownfield intent becomes structured YAML with no rewrite"],
        ["Provision an approved change", "specialized or run_*_workflow_manager tool", "The AI uses a flat contract, not raw nested Ansible data"],
        ["Watch the job run", "taskId + progress + /tasks/get", "Long-running changes stay observable and auditable"],
        ["Validate outcome", "assurance or export tool", "The same platform is used to confirm actual state"],
    ]
    add_table_slide(slide, len(rows), 3, rows, [2.8, 3.1, 6.4])


def add_use_cases_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, 4 Columns"))
    set_title(slide, "High-Value Customer Use Cases")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(
        placeholders[0],
        ["Start where Catalyst Center already owns the truth and the Ansible collection already exposes the workflow."],
        font_size=16,
    )
    fill_placeholder_lines(
        placeholders[1],
        ["Brownfield export", "Pull site, inventory, network settings, or wireless design into reusable data files."],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[2],
        ["Campus rollout", "Provision buildings, floors, fabric devices, templates, and baseline settings in a controlled path."],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[3],
        ["Day-2 changes", "Apply approved updates through typed tools instead of ad hoc API calls or one-off scripts."],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[4],
        ["Validation", "Use assurance-driven operations and exports to verify that the network matches intent after change."],
        font_size=14,
        first_line_bold=True,
    )


def add_adoption_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, Subtitle, 3 Columns"))
    set_title(slide, "A Realistic 30 / 60 / 90 Day Adoption Path")
    placeholders = content_placeholders(slide)
    fill_placeholder_lines(
        placeholders[0],
        ["Adoption works best when the first win is read-only, the second win is controlled provisioning, and the third win adds validation."],
        font_size=16,
    )
    fill_placeholder_lines(
        placeholders[1],
        [
            "First 30 days",
            "Deploy the server behind HTTPS",
            "Enable tenant-scoped credentials",
            "Start with config generators and task polling",
        ],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[2],
        [
            "Day 31 to 60",
            "Add one or two approved write paths",
            "Train operators on prompt patterns and approval workflow",
            "Capture artifacts and change evidence",
        ],
        font_size=14,
        first_line_bold=True,
    )
    fill_placeholder_lines(
        placeholders[3],
        [
            "Day 61 to 90",
            "Expand to validation and closed-loop use cases",
            "Add more workflow_manager coverage",
            "Measure adoption by task success and time saved",
        ],
        font_size=14,
        first_line_bold=True,
    )


def add_takeaways_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Title, 1 Column with Bullets"))
    set_title(slide, "Key Takeaways")
    fill_placeholder_lines(
        content_placeholders(slide)[0],
        [
            "MCP does not replace Catalyst Center automation. It gives your existing automation a safe, typed, AI-consumable front door.",
            "For Catalyst Center teams, the cleanest starting point is read-only export with playbook_config_generator modules.",
            "The same Ansible collection can then be used to provision approved changes through workflow_manager modules.",
            "Task IDs, progress events, HTTPS, Redis state, and tenant-aware access are what make the model production-worthy.",
            "Start small, prove one planning flow and one provisioning flow, then expand with validation and governance.",
        ],
        font_size=18,
    )
    add_card(
        slide,
        8.15,
        4.3,
        4.3,
        1.3,
        "Recommended first workshop",
        ["Pick one read-only export and one approved provisioning workflow that already exists in your Catalyst Center practice."],
        fill_color=SKY,
        line_color=SKY,
    )


def add_thank_you_slide(prs: Presentation):
    slide = prs.slides.add_slide(get_layout(prs, "Thank you 1"))
    set_title(slide, "Questions?")
    add_textbox(
        slide,
        0.35,
        4.25,
        7.0,
        0.8,
        ["Catalyst Center teams do not need to become AI experts first.", "They need one trusted MCP server and one useful workflow to start."],
        font_size=18,
        bold_first=True,
        color=DARK,
    )
    add_card(
        slide,
        8.15,
        2.7,
        4.2,
        1.5,
        "Leave-behind value",
        ["A customer-ready architecture, a working MCP server, and a reproducible PowerPoint deck built on the Cisco template."],
        fill_color=LIGHT,
    )


def build_presentation(template: Path, output: Path):
    converted_template, temp_dir = convert_template_path(template)
    try:
        prs = Presentation(str(converted_template))
        remove_all_slides(prs)
        add_title_slide(prs)
        add_agenda_slide(prs)
        add_segue(prs, "Why This Matters")
        add_existing_assets_slide(prs)
        add_mcp_plain_english_slide(prs)
        add_architecture_slide(prs)
        add_segue(prs, "Plan")
        add_plan_slide(prs)
        add_prompt_patterns_slide(prs)
        add_segue(prs, "Provision")
        add_tool_model_slide(prs)
        add_ppv_flow_slide(prs)
        add_segue(prs, "Validate and Operate")
        add_task_lifecycle_slide(prs)
        add_security_slide(prs)
        add_deployment_slide(prs)
        add_demo_storyboard_slide(prs)
        add_use_cases_slide(prs)
        add_adoption_slide(prs)
        add_takeaways_slide(prs)
        add_thank_you_slide(prs)
        prs.save(str(output))
    finally:
        if temp_dir is not None:
            temp_dir.cleanup()


def parse_args():
    parser = argparse.ArgumentParser(description="Build the Catalyst Center MCP customer training deck.")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE, help="Path to the Cisco PowerPoint template (.potx or .pptx).")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT, help="Output presentation path.")
    return parser.parse_args()


def main():
    args = parse_args()
    if not args.template.exists():
        raise SystemExit(f"Template not found: {args.template}")
    build_presentation(args.template, args.output)
    print(f"Wrote {args.output}")


if __name__ == "__main__":
    main()
